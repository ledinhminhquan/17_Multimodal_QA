"""Generate the submission slides.pptx (python-pptx) - ~12 concise slides for the mmqa (VQA)
system. Degrades to a Markdown outline if python-pptx is unavailable.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from ..config import AppConfig, run_dir
from ..logging_utils import get_logger
from . import charts as charts_mod
from .artifact_loader import agent_metric, base_model, load_artifacts, model_system_name, sys_accuracy

logger = get_logger(__name__)


def _pct(v):
    return f"{v*100:.1f}%" if isinstance(v, (int, float)) and not isinstance(v, bool) else "?"


def _slides(cfg: AppConfig, arts: Dict[str, Any]) -> List[Tuple[str, List[str]]]:
    mname = model_system_name(arts) or "model"
    macc = sys_accuracy(arts, mname)
    bacc = sys_accuracy(arts, "blind_prior")
    cov = agent_metric(arts, "coverage")
    res = (f"model VQA accuracy {_pct(macc)} vs blind prior {_pct(bacc)}"
           if macc is not None else "train + evaluate to populate results")
    covline = (f"agent coverage {_pct(cov)} (abstains on the rest)" if cov is not None
               else "coverage / accuracy-on-answered trade-off")
    return [
        ("Multimodal Question Answering (VQA)",
         [f"{cfg.author} - Student {cfg.student_id}", "NLP in Industry - Final Assignment",
          "Answer a natural-language question about an image",
          "Trainable VQA core (ViLT) + a 5-decision agent",
          "Abstains when unsure; never asserts certainty"]),
        ("Business Problem & Motivation",
         ["Accessibility (blind users), e-commerce Q&A, image search, triage",
          "VQA models are OVERCONFIDENT and exploit LANGUAGE PRIORS",
          "The fix = a trainable VQA core + calibrated abstention + answer constraints",
          "Only the VQA model is trained; routing + gates are algorithmic"]),
        ("Proposed Solution",
         ["Classify the question type (yes/no, count, colour, ...)",
          "Run the VQA model -> top-k answers + confidence",
          "Abstain ('unsure') when the model is uncertain",
          "Constrain the answer to the question type"]),
        ("System Architecture",
         ["ingest (D1 input gate) -> classify (D2 question type)",
          "-> answer (D3 run VQA top-k) -> calibrate (D4 abstain)",
          "-> constrain (D5 type-consistency)",
          "Runs fully offline (SceneStubVQA + synthetic scenes) for tests/CI"]),
        ("Data (VQAv2 + Synthetic Scenes)",
         ["Train: HuggingFaceM4/VQAv2 (10-annotator answers, trust_remote_code)",
          "Eval: lmms-lab/VQAv2 validation (CC-BY-4.0, clean parquet)",
          "Answer vocab: 3129 labels from the ViLT config",
          "Offline: a synthetic scene generator (shapes + embedded spec + QA)"]),
        ("The Trainable VQA Core (ViLT)",
         ["dandelin/vilt-b32-finetuned-vqa (Apache, ~113M, classification over 3129 answers)",
          "Soft-target BCE: each annotator answer gets min(1, count/3)",
          "HF Trainer, VQA-accuracy selection, bf16/tf32, early stopping",
          "Alts: BLIP (BSD, generative), GIT (MIT); H100: BLIP-2-FLAN-T5-XL (MIT)"]),
        ("The 5-Decision Agent",
         ["D1 input gate - D2 question-type classification",
          "D3 run VQA (top-k + confidence/margin/entropy)",
          "D4 CALIBRATED ABSTENTION (overconfidence is the risk)",
          "D5 TYPE-CONSISTENCY constraint (yes/no, number, colour)"]),
        ("Evaluation Results",
         [res, covline,
          "Official VQA accuracy (soft, 10-annotator) + per-answer-type",
          "Blind-prior baseline exposes the language bias; abstention buys precision"]),
        ("Deployment Overview",
         ["FastAPI POST /ask (image+question, multipart) + /ask-scene (JSON) + /healthz",
          "Gradio demo (upload image + type a question)",
          "Docker (libGL) + HF Space; lazy deps + scene-stub offline fallback",
          "Metadata-only job logging"]),
        ("Continual Learning, Monitoring & Ethics",
         ["Collect corrected/flagged answers -> re-fine-tune the VQA core",
          "monitor-log: abstain rate + confidence/answer-distribution drift + latency",
          "Privacy: user photos = sensitive (faces, homes, assistive use) -> local, no retention",
          "Over-trust risk -> abstain + flag; report blind-prior + per-type to surface bias"]),
        ("Key Takeaways & Future Work",
         ["A calibrated, type-aware, debuggable VQA pipeline",
          "Abstention + answer constraints beat raw argmax (and are safer)",
          "Future: generative BLIP-2, open-vocab answers, knowledge VQA (OK-VQA), OCR-VQA",
          "Future: better calibration (temperature scaling), unanswerable detection (VizWiz)"]),
    ]


def generate_slides(cfg: AppConfig, title: Optional[str] = None, author: Optional[str] = None,
                    out_path: Optional[str] = None) -> str:
    arts = load_artifacts(cfg)
    out_path = Path(out_path) if out_path else run_dir() / "report" / "slides.pptx"
    out_path.parent.mkdir(parents=True, exist_ok=True)
    slides = _slides(cfg, arts)
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt
    except Exception as exc:
        logger.warning("python-pptx unavailable (%s); writing markdown outline", exc)
        md = "\n\n".join(f"## {t}\n" + "\n".join(f"- {b}" for b in bs) for t, bs in slides)
        alt = out_path.with_suffix(".md")
        alt.write_text(md, encoding="utf-8")
        return str(alt)

    try:
        chart = charts_mod.accuracy_chart(arts, run_dir() / "report" / "slide_accuracy.png")
    except Exception:
        chart = None
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    accent = RGBColor(0x2B, 0x6C, 0xB0)
    for i, (t, bullets) in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
        bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
        tf = bar.text_frame; tf.text = t
        tf.paragraphs[0].font.size = Pt(28); tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        body = slide.shapes.add_textbox(Inches(0.6), Inches(1.5),
                                        Inches(8.3 if (i == 7 and chart) else 12), Inches(5.4))
        bt = body.text_frame; bt.word_wrap = True
        for j, bp in enumerate(bullets):
            p = bt.paragraphs[0] if j == 0 else bt.add_paragraph()
            p.text = "-  " + bp; p.font.size = Pt(20); p.space_after = Pt(10)
        if i == 7 and chart:
            slide.shapes.add_picture(str(chart), Inches(8.9), Inches(1.7), width=Inches(4.0))
        foot = slide.shapes.add_textbox(Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.4))
        foot.text_frame.text = f"{title or cfg.project_title} - {author or cfg.author} ({cfg.student_id})"
        foot.text_frame.paragraphs[0].font.size = Pt(9)
    prs.save(str(out_path))
    logger.info("Slides -> %s", out_path)
    return str(out_path)


__all__ = ["generate_slides"]
